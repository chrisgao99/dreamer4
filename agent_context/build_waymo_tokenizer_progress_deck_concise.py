from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


OUT = Path("/p/yufeng/tri30/agent_context/waymo_vector_tokenizer_progress.pptx")
OUT_COPY = Path("/p/yufeng/tri30/agent_context/waymo_vector_tokenizer_progress_concise.pptx")
DATA_DIR = Path("/p/yufeng/tri30/dreamer4/waymo/data/waymo_vector_dataset")
POSTER = DATA_DIR / "3e55e88e46dac74e_preview.png"
VIDEO = DATA_DIR / "3e55e88e46dac74e.mp4"

W = Inches(10.0)
H = Inches(5.625)

BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(90, 90, 90)
LIGHT = RGBColor(225, 225, 225)
BLUE = RGBColor(56, 112, 177)
GREEN = RGBColor(75, 145, 95)
RED = RGBColor(190, 82, 65)
PURPLE = RGBColor(115, 92, 160)


def font(run, size=18, bold=False, color=BLACK):
    run.font.name = "Helvetica Neue"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text(slide, body, x, y, w, h, size=18, bold=False, color=BLACK, align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    r = p.add_run()
    r.text = body
    font(r, size=size, bold=bold, color=color)
    return box


def title(slide, body):
    text(slide, body, Inches(0.17), Inches(0.11), Inches(9.32), Inches(0.63), size=26, bold=False)


def footer(slide, num):
    text(slide, str(num), Inches(9.35), Inches(5.39), Inches(0.45), Inches(0.16), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def bullets(slide, items, x=0.45, y=1.0, w=8.8, size=18, line=0.48):
    for i, item in enumerate(items):
        text(slide, f"- {item}", Inches(x), Inches(y + i * line), Inches(w), Inches(0.35), size=size)


def small_label(slide, body, x, y, w, h, color=BLACK):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.2)
    shape.adjustments[0] = 0.05
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = body
    font(r, size=13, color=color)
    return shape


def arrow(slide, x1, y1, x2, y2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = BLACK
    conn.line.width = Pt(1.0)
    conn.line.end_arrowhead = True
    return conn


def add_table(slide, rows, x, y, w, h, widths):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for i, frac in enumerate(widths):
        tbl.columns[i].width = int(w * frac)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Helvetica Neue"
                p.font.size = Pt(12 if r else 13)
                p.font.bold = bool(r == 0)
                p.font.color.rgb = BLACK
    return tbl


def main():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # 1
    s = prs.slides.add_slide(blank)
    text(s, "Waymo vector tokenizer progress", Inches(0.34), Inches(1.45), Inches(9.0), Inches(0.7), size=30)
    text(s, "Yufeng Gao\nTRI-UVA sync", Inches(0.34), Inches(2.65), Inches(5.0), Inches(0.7), size=18)
    footer(s, 1)

    # 2
    s = prs.slides.add_slide(blank)
    title(s, "What I finished last week")
    bullets(s, [
        "Waymo tf.Example -> vector NPZ filter",
        "Agent / map / traffic-light tensors with masks",
        "Ego-centric normalization and local map crop",
        "MP4 visualization for checking filtered scenes",
        "First Dreamer-style encoder smoke test",
    ], size=20, line=0.55)
    footer(s, 2)

    # 3
    s = prs.slides.add_slide(blank)
    title(s, "Waymo input after preprocessing")
    rows = [
        ["input", "shape", "features"],
        ["agents", "(K, 91, 8)", "x, y, speed, vx, vy, valid, yaw, type"],
        ["map", "(M, P, 6)", "x, y, dir_x, dir_y, type, valid"],
        ["lights", "(91, 16, 4)", "x, y, state, valid"],
        ["masks", "agent / map / light", "padding and missing values"],
    ]
    add_table(s, rows, Inches(0.4), Inches(1.1), Inches(9.1), Inches(2.3), [0.2, 0.22, 0.58])
    text(s, "Current debug setting: K=16, M=256, P=20", Inches(0.45), Inches(4.05), Inches(8.8), Inches(0.35), size=18)
    footer(s, 3)

    # 4
    s = prs.slides.add_slide(blank)
    title(s, "Filtering")
    bullets(s, [
        "ego / SDC is slot 0",
        "select closest K-1 agents",
        "crop roadgraph near ego trajectory",
        "group roadgraph samples into polylines",
        "normalize all geometry to ego current frame",
    ], x=0.35, y=0.95, w=4.15, size=18, line=0.42)
    s.shapes.add_picture(str(POSTER), Inches(4.85), Inches(0.95), width=Inches(4.75))
    text(s, "Example MP4: 91 frames, 10 Hz", Inches(4.9), Inches(4.95), Inches(4.5), Inches(0.25), size=12, color=GRAY)
    try:
        s.shapes.add_movie(str(VIDEO), Inches(8.85), Inches(4.55), Inches(0.5), Inches(0.28), poster_frame_image=str(POSTER), mime_type="video/mp4")
    except Exception:
        pass
    footer(s, 4)

    # 5
    s = prs.slides.add_slide(blank)
    title(s, "Encoder tokens")
    x0 = Inches(0.55)
    y = Inches(1.75)
    boxes = [
        ("latent\nz_t", PURPLE),
        ("agent\nK tracks", GREEN),
        ("map\nM polylines", BLUE),
        ("light\n16 slots", RED),
    ]
    for i, (label, color) in enumerate(boxes):
        x = x0 + Inches(i * 1.85)
        small_label(s, label, x, y, Inches(1.25), Inches(0.72), color)
        if i < len(boxes) - 1:
            arrow(s, x + Inches(1.25), y + Inches(0.36), x + Inches(1.68), y + Inches(0.36))
    text(s, "[latent, agent, map, light] at each timestep", Inches(0.55), Inches(3.05), Inches(9.0), Inches(0.35), size=20, align=PP_ALIGN.CENTER)
    bullets(s, [
        "latent tokens: tokenizer bottleneck",
        "agent tokens: per-agent dynamics representation",
        "scene/task/policy-agent tokens only during policy finetuning",
    ], x=0.9, y=4.0, w=8.4, size=16, line=0.34)
    footer(s, 5)

    # 6
    s = prs.slides.add_slide(blank)
    title(s, "Encoder design: reuse existing ideas")
    small_label(s, "MTR", Inches(0.8), Inches(1.25), Inches(1.35), Inches(0.55), BLUE)
    bullets(s, [
        "ego / center-object coordinates",
        "map as polylines",
        "PointNet-style polyline encoder",
    ], x=2.35, y=1.24, w=6.9, size=18, line=0.42)
    small_label(s, "Dreamer", Inches(0.8), Inches(3.15), Inches(1.35), Inches(0.55), PURPLE)
    bullets(s, [
        "space attention within timestep",
        "causal time attention across timesteps",
        "latent bottleneck tokens",
    ], x=2.35, y=3.14, w=6.9, size=18, line=0.42)
    footer(s, 6)

    # 7
    s = prs.slides.add_slide(blank)
    title(s, "Map input plan")
    text(s, "baseline now", Inches(0.7), Inches(1.2), Inches(3.7), Inches(0.35), size=20)
    bullets(s, [
        "encode all cropped map polylines",
        "repeat map tokens over time",
        "simple and small compared to image tokens",
    ], x=0.85, y=1.75, w=4.0, size=17, line=0.42)
    text(s, "next comparison", Inches(5.45), Inches(1.2), Inches(3.7), Inches(0.35), size=20)
    bullets(s, [
        "encode map separately",
        "treat map as static memory",
        "agents query map context",
    ], x=5.6, y=1.75, w=3.9, size=17, line=0.42)
    arrow(s, Inches(4.55), Inches(2.6), Inches(5.15), Inches(2.6))
    text(s, "compare reconstruction performance", Inches(1.9), Inches(4.4), Inches(6.6), Inches(0.35), size=20, align=PP_ALIGN.CENTER)
    footer(s, 7)

    # 8
    s = prs.slides.add_slide(blank)
    title(s, "Decoder plan: specific choices")
    small_label(s, "z_t", Inches(0.8), Inches(1.65), Inches(1.1), Inches(0.55), PURPLE)
    arrow(s, Inches(1.9), Inches(1.93), Inches(2.6), Inches(1.93))
    small_label(s, "decoder\nqueries", Inches(2.6), Inches(1.5), Inches(1.55), Inches(0.85), BLUE)
    arrow(s, Inches(4.15), Inches(1.93), Inches(4.85), Inches(1.93))
    small_label(s, "heads", Inches(4.85), Inches(1.65), Inches(1.1), Inches(0.55), GREEN)
    bullets(s, [
        "z-only decoder first: no encoder tokens bypass the bottleneck",
        "learned agent / light / optional map queries",
        "heads: agent state, agent valid, light state, light valid",
        "compare strict z-only vs z + static map memory",
    ], x=0.9, y=3.0, w=8.5, size=17, line=0.38)
    footer(s, 8)

    # 9
    s = prs.slides.add_slide(blank)
    title(s, "Next plan: staged training")
    text(s, "1. tokenizer first", Inches(0.55), Inches(1.05), Inches(2.5), Inches(0.35), size=20)
    bullets(s, [
        "verify z reconstructs agents / lights",
        "choose map handling baseline",
    ], x=0.75, y=1.55, w=2.65, size=15, line=0.34)
    text(s, "2. Dreamer dynamics", Inches(3.75), Inches(1.05), Inches(2.7), Inches(0.35), size=20)
    bullets(s, [
        "train z_t -> z_{t+1}",
        "roll out future z",
        "decode future agents / lights",
    ], x=3.95, y=1.55, w=2.75, size=15, line=0.34)
    text(s, "3. representation", Inches(7.0), Inches(1.05), Inches(2.4), Inches(0.35), size=20)
    bullets(s, [
        "downstream scene/task/policy-agent tokens for policy",
        "selected-agent tokens for agent-level dynamics",
    ], x=7.2, y=1.55, w=2.45, size=15, line=0.34)
    text(s, "Why tokenizer first: it is the measurable bottleneck before training the full world model.", Inches(0.7), Inches(4.25), Inches(8.8), Inches(0.35), size=18, align=PP_ALIGN.CENTER)
    footer(s, 9)

    # 10
    s = prs.slides.add_slide(blank)
    title(s, "Interaction representation evaluation")
    bullets(s, [
        "Compare learned tokens by future prediction and reconstruction",
        "Probe latent / agent / downstream task tokens for maneuver labels",
        "Retrieval: similar interaction dynamics should be nearest neighbors",
        "Optional contrastive loss after tokenizer baseline is stable",
    ], size=19, line=0.5)
    text(s, "Contrastive positives: same scenario crops, interaction-preserving augmentation, or same filtered maneuver class.", Inches(0.55), Inches(4.35), Inches(8.9), Inches(0.35), size=15, color=GRAY)
    footer(s, 10)

    prs.save(OUT)
    prs.save(OUT_COPY)
    print(OUT)
    print(OUT_COPY)


if __name__ == "__main__":
    main()
