from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


OUT = Path("/p/yufeng/tri30/agent_context/waymo_vector_tokenizer_progress.pptx")
DATA_DIR = Path("/p/yufeng/tri30/dreamer4/waymo/data/waymo_vector_dataset")
VIDEO = DATA_DIR / "3e55e88e46dac74e.mp4"
POSTER = DATA_DIR / "3e55e88e46dac74e_preview.png"


WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

COLORS = {
    "ink": RGBColor(26, 32, 44),
    "muted": RGBColor(88, 101, 116),
    "line": RGBColor(213, 220, 228),
    "bg": RGBColor(247, 249, 252),
    "panel": RGBColor(255, 255, 255),
    "navy": RGBColor(24, 63, 88),
    "teal": RGBColor(32, 116, 111),
    "green": RGBColor(60, 142, 94),
    "rust": RGBColor(182, 91, 59),
    "blue": RGBColor(59, 105, 180),
    "gold": RGBColor(207, 159, 65),
    "purple": RGBColor(111, 88, 156),
}


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=COLORS["line"], width=1):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_bg(slide):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, WIDE_W, WIDE_H)
    set_fill(rect, COLORS["bg"])
    rect.line.fill.background()
    slide.shapes._spTree.remove(rect._element)
    slide.shapes._spTree.insert(2, rect._element)


def textbox(slide, text, x, y, w, h, size=18, bold=False, color=COLORS["ink"], align=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Aptos"
    return box


def title(slide, text, subtitle=None):
    textbox(slide, text, Inches(0.58), Inches(0.38), Inches(11.7), Inches(0.48), size=24, bold=True)
    if subtitle:
        textbox(slide, subtitle, Inches(0.6), Inches(0.88), Inches(11.8), Inches(0.33), size=10, color=COLORS["muted"])
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.6), Inches(1.27), Inches(12.1), Inches(0.02))
    set_fill(line, COLORS["line"])
    line.line.fill.background()


def bullet_list(slide, items, x, y, w, h, size=15, color=COLORS["ink"], gap=0.98):
    top = y
    for item in items:
        dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x, top + Inches(0.09), Inches(0.08), Inches(0.08))
        set_fill(dot, COLORS["teal"])
        dot.line.fill.background()
        textbox(slide, item, x + Inches(0.18), top, w - Inches(0.18), Inches(0.45), size=size, color=color)
        top += Inches(gap * 0.36)


def card(slide, x, y, w, h, heading, body, accent=COLORS["teal"]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shape, COLORS["panel"])
    set_line(shape)
    shape.adjustments[0] = 0.08
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    set_fill(bar, accent)
    bar.line.fill.background()
    textbox(slide, heading, x + Inches(0.25), y + Inches(0.18), w - Inches(0.45), Inches(0.28), size=14, bold=True)
    bullet_list(slide, body, x + Inches(0.27), y + Inches(0.65), w - Inches(0.45), h - Inches(0.8), size=11, gap=0.86)
    return shape


def table(slide, rows, x, y, w, h, col_fracs, header=True):
    tbl = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h).table
    for i, frac in enumerate(col_fracs):
        tbl.columns[i].width = int(w * frac)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            fill = COLORS["navy"] if (header and r == 0) else COLORS["panel"]
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(10 if r else 10.5)
                p.font.bold = bool(header and r == 0)
                p.font.color.rgb = RGBColor(255, 255, 255) if (header and r == 0) else COLORS["ink"]
    return tbl


def box(slide, text, x, y, w, h, fill, font=12, bold=True):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    set_fill(shape, fill)
    set_line(shape, fill, width=0.5)
    shape.adjustments[0] = 0.08
    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(font)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(255, 255, 255)
    return shape


def arrow(slide, x1, y1, x2, y2, color=COLORS["muted"]):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(1.5)
    conn.line.end_arrowhead = True
    return conn


def add_footer(slide, idx):
    textbox(slide, f"{idx}", Inches(12.45), Inches(7.08), Inches(0.35), Inches(0.2), size=9, color=COLORS["muted"], align=PP_ALIGN.RIGHT)
    textbox(slide, "Waymo vector tokenizer progress", Inches(0.62), Inches(7.08), Inches(4.0), Inches(0.2), size=9, color=COLORS["muted"])


def main():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    # 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    textbox(slide, "Waymo Vector Tokenizer Progress", Inches(0.72), Inches(0.78), Inches(11.8), Inches(0.7), size=34, bold=True, color=COLORS["navy"])
    textbox(slide, "Dreamer-4-style world model input pipeline and encoder prototype", Inches(0.75), Inches(1.55), Inches(11.2), Inches(0.35), size=17, color=COLORS["muted"])
    card(slide, Inches(0.78), Inches(2.35), Inches(3.7), Inches(2.75), "Built last week", [
        "Waymo tf.Example -> fixed vector tensors",
        "Local map crop and polyline chunks",
        "Traffic-light sequence tensors",
        "MP4/PNG visualization for inspection",
        "Encoder smoke tests on 32-step and 11-step windows",
    ], COLORS["green"])
    card(slide, Inches(4.85), Inches(2.35), Inches(3.7), Inches(2.75), "Current model state", [
        "Implemented encoder only",
        "Token layout: latent, agent, map, light",
        "MTR-style map polyline stem",
        "Dreamer-style block-causal space-time attention",
        "Decoder and training losses are next",
    ], COLORS["blue"])
    card(slide, Inches(8.92), Inches(2.35), Inches(3.7), Inches(2.75), "Boss-facing message", [
        "Not rebuilding the full stack",
        "Reuse MTR ideas for vector/map geometry",
        "Reuse Dreamer ideas for token/time modeling",
        "Add scene/task/policy-agent tokens only during policy finetuning",
        "Compare map handling baselines before optimizing",
    ], COLORS["rust"])
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Progress Snapshot", "Files under /p/yufeng/tri30/dreamer4/waymo")
    rows = [
        ["Component", "What it does", "Status"],
        ["waymo_vector_filter.py", "Parses TFRecords without TensorFlow and saves fixed NPZ tensors", "Implemented"],
        ["waymo_vector_dataset.py", "Loads NPZ files as PyTorch tensors for model input", "Implemented"],
        ["visualize_waymo_vector_npz.py", "Draws map, agents, heading arrows, trails, and lights to MP4/PNG", "Implemented"],
        ["vector_tokenizer_encoder.py", "Runs first Dreamer-4-style vector encoder", "Encoder smoke-tested"],
        ["Decoder + train script", "Reconstruct/predict agents, lights, optional cropped map", "Planned next"],
    ]
    table(slide, rows, Inches(0.7), Inches(1.65), Inches(11.95), Inches(4.1), [0.25, 0.55, 0.2])
    textbox(slide, "Debug dataset: 10 filtered NPZ scenarios; two MP4 examples at 91 frames / 10 Hz.", Inches(0.75), Inches(6.1), Inches(11.4), Inches(0.35), size=14, bold=True, color=COLORS["teal"])
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Preprocess Waymo Input", "Raw Waymo motion examples have 10 past + 1 current + 80 future = 91 states.")
    rows = [
        ["Tensor", "Shape", "Feature order"],
        ["agents", "(K, 91, 8)", "x, y, speed, vx, vy, valid, yaw, type"],
        ["agent_mask", "(K,)", "valid selected slot"],
        ["map_polylines", "(M, P, 6)", "x, y, dir_x, dir_y, type, valid"],
        ["map_mask", "(M, P)", "valid roadgraph point"],
        ["lights", "(91, 16, 4)", "x, y, state, valid"],
        ["light_mask", "(91, 16)", "valid traffic light slot"],
        ["ego_origin_xy, ego_heading", "(2,), scalar", "saved for converting predictions back to world frame"],
    ]
    table(slide, rows, Inches(0.7), Inches(1.55), Inches(12.0), Inches(4.8), [0.24, 0.18, 0.58])
    textbox(slide, "Current debug settings: K=16, M=256, P=20. The encoder also accepts variable time windows such as T=32 and T=11.", Inches(0.75), Inches(6.55), Inches(11.6), Inches(0.35), size=13, color=COLORS["muted"])
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Filtering And Coordinate Frame", "The goal is fixed tensors that keep local interaction context without feeding the whole Waymo roadgraph.")
    card(slide, Inches(0.72), Inches(1.65), Inches(3.65), Inches(4.55), "Agent filtering", [
        "Ego/SDC is always slot 0",
        "Select closest K-1 agents by minimum valid distance to ego",
        "Pad missing slots and keep masks",
        "Current option can select over all 91 states; history-only selection is available for causal deployment",
    ], COLORS["green"])
    card(slide, Inches(4.82), Inches(1.65), Inches(3.65), Inches(4.55), "Map filtering", [
        "Start from 30,000 roadgraph samples",
        "Keep map IDs with points near the ego trajectory",
        "Split kept IDs into fixed-length polyline chunks",
        "Cap and pad to M polylines x P points",
    ], COLORS["blue"])
    card(slide, Inches(8.92), Inches(1.65), Inches(3.65), Inches(4.55), "Ego-centric normalization", [
        "Use ego current pose as origin",
        "Rotate positions, velocities, map directions, and light positions into ego frame",
        "Keep semantic fields unchanged",
        "This follows MTR center-object coordinates; our center object is the ego/SDC",
    ], COLORS["rust"])
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Example Waymo Visualization", "Generated from filtered NPZ to verify what the model actually receives.")
    pic = slide.shapes.add_picture(str(POSTER), Inches(0.72), Inches(1.55), width=Inches(7.6))
    set_line(pic, COLORS["line"])
    try:
        movie = slide.shapes.add_movie(str(VIDEO), Inches(8.7), Inches(1.65), Inches(3.6), Inches(2.02), poster_frame_image=str(POSTER), mime_type="video/mp4")
        set_line(movie, COLORS["line"])
    except Exception:
        textbox(slide, "MP4 file is packaged next to the deck:\n" + str(VIDEO), Inches(8.7), Inches(1.65), Inches(3.6), Inches(1.0), size=11, color=COLORS["muted"])
    card(slide, Inches(8.55), Inches(4.05), Inches(3.8), Inches(1.9), "What is shown", [
        "Gray: cropped map polylines",
        "Green: ego agent",
        "Other colors: selected nearby agents",
        "Arrows/trails: heading and recent motion",
        "Lights: red/yellow/green when valid",
    ], COLORS["teal"])
    textbox(slide, "Example MP4: /p/yufeng/tri30/dreamer4/waymo/data/waymo_vector_dataset/3e55e88e46dac74e.mp4", Inches(0.78), Inches(6.62), Inches(11.7), Inches(0.28), size=10, color=COLORS["muted"])
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Encoder Token Design", "At every timestep, the encoder builds one structured token set.")
    y = Inches(2.05)
    x0 = Inches(0.8)
    specs = [
        ("latent", "bottleneck z_t", COLORS["purple"]),
        ("agent", "K selected tracks", COLORS["green"]),
        ("map", "M local polylines", COLORS["blue"]),
        ("light", "L traffic lights", COLORS["rust"]),
    ]
    xs = [x0, Inches(3.0), Inches(5.2), Inches(7.45)]
    for i, (name, desc, col) in enumerate(specs):
        box(slide, name + "\n" + desc, xs[i], y, Inches(1.72), Inches(0.9), col, font=12)
        if i < len(specs) - 1:
            arrow(slide, xs[i] + Inches(1.72), y + Inches(0.45), xs[i + 1] - Inches(0.12), y + Inches(0.45))
    textbox(slide, "[latent_1..N, agent_1..K, map_1..M, light_1..L]", Inches(1.2), Inches(3.25), Inches(10.8), Inches(0.35), size=20, bold=True, color=COLORS["ink"], align=PP_ALIGN.CENTER)
    card(slide, Inches(0.8), Inches(4.2), Inches(3.55), Inches(1.65), "No policy tokens", [
        "Tokenizer pretraining follows Dreamer 4 causal-tokenizer phase",
        "Scene/task/policy-agent tokens are inserted only during policy finetuning",
    ], COLORS["navy"])
    card(slide, Inches(4.85), Inches(4.2), Inches(3.55), Inches(1.65), "Latent tokens", [
        "Compact tokenizer bottleneck",
        "Output z has shape (B, T, N_latents, D_bottleneck)",
    ], COLORS["purple"])
    card(slide, Inches(8.9), Inches(4.2), Inches(3.55), Inches(1.65), "Agent tokens", [
        "Per-agent dynamics representations",
        "Can support downstream behavior/interaction analysis",
    ], COLORS["green"])
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Encoder: Reuse MTR And Dreamer Ideas", "This is a thin adaptation, not a full rebuild of existing work.")
    card(slide, Inches(0.72), Inches(1.6), Inches(3.75), Inches(4.7), "From MTR", [
        "Ego/center-object coordinate frame",
        "Polyline map representation",
        "PointNet-style polyline encoder: point MLP, masked max-pool, local/global fusion",
        "Agent-map interaction motivation",
    ], COLORS["blue"])
    card(slide, Inches(4.82), Inches(1.6), Inches(3.75), Inches(4.7), "From Dreamer 4", [
        "Tokenized world-model view",
        "Space attention within each timestep",
        "Causal time attention across persistent slots",
        "Bottleneck latent tokens for dynamics training",
    ], COLORS["purple"])
    card(slide, Inches(8.92), Inches(1.6), Inches(3.75), Inches(4.7), "Repo-native additions", [
        "Simple agent and traffic-light MLP stems",
        "Scene/task/policy-agent tokens reserved for policy finetuning",
        "Waymo-specific filters, masks, and visualization",
        "Future decoder/training loop integrated with Dreamer code",
    ], COLORS["rust"])
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Map Handling Baseline And Planned Comparison", "The current encoder intentionally starts with the simplest strong baseline.")
    card(slide, Inches(0.72), Inches(1.55), Inches(5.65), Inches(4.9), "Current baseline: temporal map tokens", [
        "Encode all cropped local map polylines once",
        "Repeat map tokens at every timestep: (B, M, D) -> (B, T, M, D)",
        "Agents, lights, and latents can directly attend to map tokens in space attention",
        "Reasonable first baseline because cropped vector map is small compared with Dreamer image-token inputs",
    ], COLORS["green"])
    card(slide, Inches(6.95), Inches(1.55), Inches(5.65), Inches(4.9), "Next method: static map memory", [
        "Encode map separately as static context",
        "Dynamic tokens query map memory through cross-attention",
        "Avoid repeated map time slots and make static/dynamic separation cleaner",
        "Compare the two methods by reconstruction quality, future prediction quality, and compute/memory",
    ], COLORS["blue"])
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Decoder And Training Plan", "Decoder is the next implementation step; map reconstruction stays optional.")
    box(slide, "z_t latents", Inches(0.8), Inches(2.0), Inches(1.75), Inches(0.75), COLORS["purple"])
    arrow(slide, Inches(2.55), Inches(2.38), Inches(3.35), Inches(2.38))
    box(slide, "decoder queries\nagent / light / optional map", Inches(3.35), Inches(1.85), Inches(3.2), Inches(1.05), COLORS["navy"], font=11)
    arrow(slide, Inches(6.55), Inches(2.38), Inches(7.4), Inches(2.38))
    box(slide, "space-time decoder", Inches(7.4), Inches(2.0), Inches(2.05), Inches(0.75), COLORS["teal"])
    arrow(slide, Inches(9.45), Inches(2.38), Inches(10.2), Inches(2.38))
    box(slide, "reconstruction\nand prediction heads", Inches(10.2), Inches(1.85), Inches(2.3), Inches(1.05), COLORS["rust"], font=11)
    card(slide, Inches(0.8), Inches(4.0), Inches(3.55), Inches(1.75), "Agent targets", [
        "x, y, speed, vx, vy, valid",
        "yaw optional",
        "masked reconstruction and future prediction",
    ], COLORS["green"])
    card(slide, Inches(4.85), Inches(4.0), Inches(3.55), Inches(1.75), "Light targets", [
        "state classification",
        "valid prediction",
        "optional x/y if useful",
    ], COLORS["gold"])
    card(slide, Inches(8.9), Inches(4.0), Inches(3.55), Inches(1.75), "Map target", [
        "optional auxiliary loss only",
        "reconstruct cropped input map, not global Waymo map",
        "main objective remains dynamic future state",
    ], COLORS["blue"])
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(blank)
    add_bg(slide)
    title(slide, "Verification And Next Steps", "Current result is a working data path plus encoder prototype.")
    rows = [
        ["Check", "Result"],
        ["Filtered data", "10 NPZ examples with agents (16,91,8), maps (256,20,6), lights (91,16,4)"],
        ["Visualization", "Two MP4 examples, each 91 frames at 10 Hz, plus preview PNG sheets"],
        ["Encoder T=32", "z (2,32,8,32), agent (2,32,16,128), map (2,32,256,128), light (2,32,16,128)"],
        ["Encoder T=11", "z (2,11,8,32), supports observed-prefix input"],
    ]
    table(slide, rows, Inches(0.75), Inches(1.55), Inches(11.9), Inches(2.85), [0.24, 0.76])
    card(slide, Inches(0.8), Inches(4.85), Inches(3.55), Inches(1.45), "1. Decoder", [
        "Implement query decoder and masked reconstruction losses",
    ], COLORS["rust"])
    card(slide, Inches(4.85), Inches(4.85), Inches(3.55), Inches(1.45), "2. Map comparison", [
        "Temporal repeated map baseline vs static map memory",
    ], COLORS["blue"])
    card(slide, Inches(8.9), Inches(4.85), Inches(3.55), Inches(1.45), "3. Scale data", [
        "Run filter over more Waymo shards and track reconstruction metrics",
    ], COLORS["green"])
    add_footer(slide, 10)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
